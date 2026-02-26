"""Generate hackathon presentation PPTX with precise layout control.

Uses python-pptx runs for bold/normal mixed text formatting.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_BLUE = RGBColor(0x00, 0x4E, 0x8C)
FONT_TITLE = "Segoe UI"
FONT_BODY = "Calibri"

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "presentations", "ghcpsdknotify_v4.pptx")
IMAGE_PATH = os.path.join(SCRIPT_DIR, "assets", "architecture.png")


def _add_title_bar(slide, text, top=Inches(0)):
    """Add a blue title bar across the top of a slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=top,
        width=SLIDE_WIDTH, height=Inches(1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)


def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _add_rich_paragraph(tf, segments, font_size=Pt(12), color=BLACK,
                        font_name=FONT_BODY, space_after=Pt(6),
                        space_before=Pt(0), level=0, first=False):
    """Add a paragraph with mixed bold/normal runs.

    segments: list of (text, bold) tuples.
    Example: [("Normal ", False), ("Bold part", True), (" normal again", False)]
    """
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_after = space_after
    p.space_before = space_before
    p.level = level

    for i, (text, bold) in enumerate(segments):
        if i == 0 and first:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name

    return p


def _add_subtitle(tf, text, first=True):
    """Add a subtitle paragraph."""
    return _add_rich_paragraph(
        tf, [(text, True)],
        font_size=Pt(16), color=BLUE, font_name=FONT_TITLE,
        space_after=Pt(10), first=first,
    )


def build_slide1(prs):
    """Slide 1: Problem & Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    _add_title_bar(slide, "Problem & Solution")

    # ── Left column: Problem ──
    tf_left = _add_textbox(
        slide,
        left=Inches(0.5), top=Inches(1.3),
        width=Inches(5.8), height=Inches(5.5),
    )
    _add_subtitle(tf_left, "Problem: The Barrier to Leveraging Gen AI", first=True)

    _add_rich_paragraph(tf_left, [
        ("Context Window Limits", True),
        (" \u2014 LLMs are powerful, but users need to ingest vast personal knowledge "
         "(notes, memos, learning records) into limited context windows. "
         "Users want ", False),
        ("personalized information, not generic answers.", True),
    ], font_size=Pt(12), space_after=Pt(10))

    _add_rich_paragraph(tf_left, [
        ("Privacy Concerns", True),
        (" \u2014 Local files often contain sensitive data or business notes. "
         "Uploading them to cloud APIs poses significant risks.", False),
    ], font_size=Pt(12), space_after=Pt(10))

    # ── Vertical divider ──
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(6.55), top=Inches(1.5),
        width=Inches(0.02), height=Inches(5.0),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    divider.line.fill.background()

    # ── Right column: Solution ──
    tf_right = _add_textbox(
        slide,
        left=Inches(6.9), top=Inches(1.3),
        width=Inches(5.8), height=Inches(5.5),
    )
    _add_subtitle(tf_right, "Solution: Personal AI Daily Briefing Agent", first=True)

    # ── Key value proposition (highlighted) ──
    _add_rich_paragraph(tf_right, [
        ("Your notes become the context \u2014 delivering ", False),
        ("personalized news, timely reminders, and spaced-repetition learning", True),
        (" tailored to your own knowledge base.", False),
    ], font_size=Pt(13), color=DARK_BLUE, space_after=Pt(14))

    # ── Features ──
    _add_rich_paragraph(tf_right, [
        ("Feature A \u2014 Daily Briefing: ", True),
        ("Searches & summarizes the latest news based on note topics "
         "(Copilot SDK + Bing + WorkIQ MCP)", False),
    ], font_size=Pt(11), space_after=Pt(8))

    _add_rich_paragraph(tf_right, [
        ("Feature B \u2014 Review Quiz: ", True),
        ("Auto-generates Q1 (multiple-choice) & Q2 (free-form) from notes, "
         "scored by LLM (Copilot SDK)", False),
    ], font_size=Pt(11), space_after=Pt(8))

    _add_rich_paragraph(tf_right, [
        ("Spaced Repetition: ", True),
        ("SM-2 algorithm auto-adjusts intervals "
         "(1 \u2192 3 \u2192 7 \u2192 14 \u2192 30 \u2192 60 days)", False),
    ], font_size=Pt(11), space_after=Pt(12))

    _add_rich_paragraph(tf_right, [
        ("Context Strategy: ", True),
        ("Weighted random + discovery rotation, up to 20 files/session", False),
    ], font_size=Pt(10), color=GRAY, space_after=Pt(4))


def build_slide2(prs):
    """Slide 2: Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    _add_title_bar(slide, "Architecture")

    # ── Architecture image (centered) ──
    if os.path.exists(IMAGE_PATH):
        from PIL import Image
        with Image.open(IMAGE_PATH) as img:
            img_w, img_h = img.size

        max_width = Inches(10)
        max_height = Inches(4.8)

        ratio = min(max_width / Emu(int(img_w * 914400 / 96)),
                    max_height / Emu(int(img_h * 914400 / 96)))
        pic_width = int(img_w * 914400 / 96 * ratio)
        pic_height = int(img_h * 914400 / 96 * ratio)

        left = int((SLIDE_WIDTH - pic_width) / 2)
        top = Inches(1.3)

        slide.shapes.add_picture(IMAGE_PATH, left, top, pic_width, pic_height)
    else:
        tf = _add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(3))
        _add_rich_paragraph(tf, [
            ("[Architecture diagram: assets/architecture.png not found]", False),
        ], font_size=Pt(18), color=GRAY, first=True)

    # ── Flow description (moved from Slide 1) ──
    tf_flow = _add_textbox(
        slide,
        left=Inches(0.5), top=Inches(6.0),
        width=Inches(12.3), height=Inches(0.4),
    )
    p = _add_rich_paragraph(tf_flow, [
        ("Windows system-tray app \u2014 scheduled jobs invoke GitHub Copilot SDK \u2192 "
         "Markdown output \u2192 Windows toast notifications \u2192 HTML viewer", False),
    ], font_size=Pt(11), color=GRAY, first=True)
    p.alignment = PP_ALIGN.CENTER

    # ── Tech Stack ──
    tf_tech = _add_textbox(
        slide,
        left=Inches(0.5), top=Inches(6.4),
        width=Inches(12.3), height=Inches(0.4),
    )
    p = _add_rich_paragraph(tf_tech, [
        ("Tech Stack:  ", True),
        ("Python 3.12  \u00b7  uv  \u00b7  Copilot SDK  \u00b7  APScheduler  "
         "\u00b7  pystray  \u00b7  winotify  \u00b7  tkinterweb  \u00b7  pytest", False),
    ], font_size=Pt(11), color=GRAY, first=True)
    p.alignment = PP_ALIGN.CENTER

    # ── GitHub / License ──
    tf_footer = _add_textbox(
        slide,
        left=Inches(0.5), top=Inches(6.8),
        width=Inches(12.3), height=Inches(0.4),
    )
    p = _add_rich_paragraph(tf_footer, [
        ("GitHub: https://github.com/shitada/ghcpsdknotify   |   License: MIT", False),
    ], font_size=Pt(11), color=BLUE, first=True)
    p.alignment = PP_ALIGN.CENTER


def main():
    import io
    tpl = os.path.join(SCRIPT_DIR, "assets", "default_template.pptx")
    with open(tpl, "rb") as f:
        stream = io.BytesIO(f.read())
    prs = Presentation(stream)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide1(prs)
    build_slide2(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    # Save via BytesIO to avoid OneDrive path issues with python-pptx 1.0.2
    out_stream = io.BytesIO()
    prs.save(out_stream)
    with open(OUTPUT_PATH, "wb") as f:
        f.write(out_stream.getvalue())
    print(f"PPTX saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
